from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parent
PHOTO_DIR = ROOT / "assets-treinamento"
OUT = ROOT / "Treinamento - Chumbamento de Passantes.pptx"

COLORS = {
    "bg": RGBColor(247, 248, 246),
    "ink": RGBColor(35, 40, 45),
    "muted": RGBColor(95, 103, 110),
    "blue": RGBColor(31, 91, 122),
    "green": RGBColor(45, 125, 78),
    "red": RGBColor(166, 55, 55),
    "amber": RGBColor(176, 124, 45),
    "panel": RGBColor(232, 235, 233),
    "white": RGBColor(255, 255, 255),
}

photos = {
    number: PHOTO_DIR / f"foto-{number}.jpg" for number in range(1, 11)
}


def set_slide_bg(slide, color=COLORS["bg"]):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text(slide, text, x, y, w, h, size=24, bold=False, color=None, align=None):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.margin_left = Inches(0.08)
    frame.margin_right = Inches(0.08)
    frame.margin_top = Inches(0.04)
    frame.margin_bottom = Inches(0.04)
    frame.vertical_anchor = MSO_ANCHOR.TOP
    p = frame.paragraphs[0]
    p.text = text
    p.font.name = "Aptos"
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color or COLORS["ink"]
    if align:
        p.alignment = align
    return box


def add_title(slide, title, subtitle=None):
    add_text(slide, title, 0.55, 0.35, 12.2, 0.45, size=26, bold=True, color=COLORS["ink"])
    if subtitle:
        add_text(slide, subtitle, 0.58, 0.83, 11.9, 0.3, size=10.5, color=COLORS["muted"])
    line = slide.shapes.add_shape(1, Inches(0.6), Inches(1.16), Inches(12.1), Inches(0.02))
    line.fill.solid()
    line.fill.fore_color.rgb = COLORS["blue"]
    line.line.color.rgb = COLORS["blue"]


def add_panel(slide, x, y, w, h, fill=COLORS["panel"], line=None):
    shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line or fill
    return shape


def add_bullets(slide, items, x, y, w, h, size=18, color=None, bullet=True):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.margin_left = Inches(0.12)
    frame.margin_right = Inches(0.12)
    frame.margin_top = Inches(0.08)
    frame.margin_bottom = Inches(0.08)
    for i, item in enumerate(items):
        p = frame.paragraphs[0] if i == 0 else frame.add_paragraph()
        p.text = item
        p.font.name = "Aptos"
        p.font.size = Pt(size)
        p.font.color.rgb = color or COLORS["ink"]
        p.space_after = Pt(7)
        if bullet:
            p.level = 0
    return box


def add_picture_fit(slide, image_path, x, y, w, h):
    with Image.open(image_path) as im:
        iw, ih = im.size
    box_ratio = w / h
    img_ratio = iw / ih
    if img_ratio > box_ratio:
        pic_w = w
        pic_h = w / img_ratio
        px = x
        py = y + (h - pic_h) / 2
    else:
        pic_h = h
        pic_w = h * img_ratio
        px = x + (w - pic_w) / 2
        py = y
    pic = slide.shapes.add_picture(str(image_path), Inches(px), Inches(py), Inches(pic_w), Inches(pic_h))
    return pic


def photo_card(slide, number, x, y, w, h, status, note, ok=True):
    add_panel(slide, x, y, w, h, fill=COLORS["white"], line=COLORS["panel"])
    add_picture_fit(slide, photos[number], x + 0.08, y + 0.08, w - 0.16, h - 0.72)
    color = COLORS["green"] if ok else COLORS["red"]
    add_text(slide, f"Foto {number} | {status}", x + 0.12, y + h - 0.56, w - 0.24, 0.2, size=10.5, bold=True, color=color)
    add_text(slide, note, x + 0.12, y + h - 0.34, w - 0.24, 0.26, size=8.5, color=COLORS["muted"])


def build():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    # 1
    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide, RGBColor(235, 238, 236))
    add_text(slide, "Treinamento", 0.75, 0.7, 5.2, 0.4, size=19, color=COLORS["blue"], bold=True)
    add_text(slide, "Chumbamento de passantes", 0.72, 1.13, 7.2, 1.05, size=34, bold=True)
    add_text(slide, "Padrão de execução para deixar o entorno dos tubos liso, nivelado e pronto para impermeabilização.", 0.77, 2.25, 5.8, 0.7, size=16, color=COLORS["muted"])
    add_picture_fit(slide, photos[1], 7.25, 0.45, 5.3, 6.55)
    add_panel(slide, 0.8, 5.9, 5.6, 0.55, fill=COLORS["blue"])
    add_text(slide, "Objetivo: reduzir retrabalho, falhas de aderência e risco de infiltração nos pontos críticos.", 0.92, 6.03, 5.35, 0.23, size=11.5, color=COLORS["white"], bold=True)

    # 2
    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide)
    add_title(slide, "Por que este serviço é crítico?", "O passante é uma interrupção na superfície: qualquer falha vira caminho preferencial para água.")
    add_panel(slide, 0.7, 1.55, 3.8, 4.85, fill=COLORS["white"])
    add_panel(slide, 4.75, 1.55, 3.8, 4.85, fill=COLORS["white"])
    add_panel(slide, 8.8, 1.55, 3.8, 4.85, fill=COLORS["white"])
    add_text(slide, "Estanqueidade", 0.96, 1.85, 3.2, 0.35, size=21, bold=True, color=COLORS["blue"], align=PP_ALIGN.CENTER)
    add_text(slide, "A impermeabilização depende de base contínua e bem aderida para impedir passagem de água no encontro piso/tubo.", 1.0, 2.45, 3.15, 1.45, size=16, color=COLORS["ink"], align=PP_ALIGN.CENTER)
    add_text(slide, "Aderência", 5.0, 1.85, 3.3, 0.35, size=21, bold=True, color=COLORS["blue"], align=PP_ALIGN.CENTER)
    add_text(slide, "Poeira, ninho, ondulação e partes soltas reduzem contato do impermeabilizante e favorecem descolamento.", 5.05, 2.45, 3.15, 1.45, size=16, color=COLORS["ink"], align=PP_ALIGN.CENTER)
    add_text(slide, "Durabilidade", 9.05, 1.85, 3.3, 0.35, size=21, bold=True, color=COLORS["blue"], align=PP_ALIGN.CENTER)
    add_text(slide, "Quando o chumbamento nasce correto, a proteção mecânica e o acabamento trabalham sem ponto alto, fissura ou acúmulo de água.", 9.1, 2.45, 3.15, 1.55, size=16, color=COLORS["ink"], align=PP_ALIGN.CENTER)
    add_text(slide, "Base técnica usada no treinamento: NBR 9575, NBR 9574, práticas de preparo de substrato e boletins de fabricantes.", 0.75, 6.75, 11.9, 0.28, size=10, color=COLORS["muted"])

    # 3
    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide)
    add_title(slide, "Padrão aceito", "Como deve ficar antes de liberar para impermeabilização.")
    add_bullets(slide, [
        "Tubo/passante rigidamente fixado, sem movimentação ao toque.",
        "Argamassa bem compactada no entorno, sem vazios, ninhos, trincas ou partes soltas.",
        "Acabamento liso, coeso e sem protuberâncias.",
        "Nível regular no entorno do passante, sem ondulações que criem ponto alto ou empoçamento.",
        "Superfície limpa, sem pó, graxa, resíduos soltos ou excesso de nata fraca.",
        "Arremate compatível com o sistema de impermeabilização indicado para a área.",
    ], 0.75, 1.55, 5.55, 4.9, size=18)
    photo_card(slide, 5, 6.75, 1.55, 2.75, 4.7, "Certo", "Base lisa e contínua junto ao tubo.", True)
    photo_card(slide, 7, 9.75, 1.55, 2.75, 4.7, "Certo", "Entorno regularizado para receber impermeabilização.", True)

    # 4
    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide)
    add_title(slide, "Exemplos corretos", "Fotos com chumbamento liso e pronto para impermeabilização, após limpeza final.")
    photo_card(slide, 1, 0.55, 1.45, 2.95, 5.2, "Certo", "Área preenchida e regularizada.", True)
    photo_card(slide, 2, 3.75, 1.45, 2.95, 5.2, "Certo", "Arremate contínuo ao redor do passante.", True)
    photo_card(slide, 3, 6.95, 1.45, 2.95, 5.2, "Certo", "Passante alinhado e base fechada.", True)
    photo_card(slide, 5, 10.15, 1.45, 2.95, 5.2, "Certo", "Superfície lisa no perímetro.", True)

    # 5
    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide)
    add_title(slide, "Mais exemplos corretos", "O acabamento deve permitir continuidade da impermeabilização no encontro com o tubo.")
    photo_card(slide, 6, 1.0, 1.45, 4.0, 5.25, "Certo", "Chumbamento regular e sem ressalto crítico.", True)
    photo_card(slide, 7, 5.25, 1.45, 4.0, 5.25, "Certo", "Regularização homogênea no entorno.", True)
    add_panel(slide, 9.7, 1.6, 2.75, 4.9, fill=COLORS["white"])
    add_text(slide, "Atenção antes de liberar", 9.95, 1.95, 2.25, 0.35, size=18, bold=True, color=COLORS["amber"], align=PP_ALIGN.CENTER)
    add_text(slide, "Mesmo nos exemplos certos, remover poeira e resíduos soltos antes da imprimação ou da primeira demão. O aspecto final esperado é superfície firme, limpa, lisa e sem ondulações.", 10.0, 2.55, 2.2, 2.15, size=14, color=COLORS["ink"], align=PP_ALIGN.CENTER)

    # 6
    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide)
    add_title(slide, "Exemplos errados", "Fotos 4, 8, 9 e 10: não liberar para impermeabilização.")
    photo_card(slide, 4, 0.55, 1.45, 2.95, 5.2, "Errado", "Desnível e acabamento ondulado.", False)
    photo_card(slide, 8, 3.75, 1.45, 2.95, 5.2, "Errado", "Vazio junto à parede e tubo sem arremate adequado.", False)
    photo_card(slide, 9, 6.95, 1.45, 2.95, 5.2, "Errado", "Irregularidades, poeira e textura fraca.", False)
    photo_card(slide, 10, 10.15, 1.45, 2.95, 5.2, "Errado", "Ondulação e borda levantada.", False)

    # 7
    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide)
    add_title(slide, "O que corrigir nos casos errados", "Critério prático para retrabalho antes da impermeabilização.")
    add_bullets(slide, [
        "Remover material solto, nata fraca, poeira e restos de concreto/argamassa.",
        "Abrir e recompor falhas até encontrar base firme e aderente.",
        "Rechumbar o passante com argamassa compatível, bem compactada e sem vazios.",
        "Regularizar o entorno no nível correto, sem lombada, rebaixo, ondulação ou ponto que acumule água.",
        "Desempenar e alisar a superfície, evitando cantos vivos e ressaltos junto ao tubo.",
        "Aguardar cura/liberação conforme procedimento da obra e ficha do produto impermeabilizante.",
    ], 0.75, 1.55, 5.75, 4.95, size=17.5)
    add_picture_fit(slide, photos[10], 7.15, 1.55, 5.0, 4.9)
    add_panel(slide, 7.15, 6.1, 5.0, 0.5, fill=COLORS["red"])
    add_text(slide, "Não corrigir com uma demão mais grossa de impermeabilizante. A base precisa ser corrigida antes.", 7.32, 6.23, 4.65, 0.18, size=10.5, bold=True, color=COLORS["white"], align=PP_ALIGN.CENTER)

    # 8
    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide)
    add_title(slide, "Passo a passo recomendado", "Sequência simples para padronizar a execução.")
    steps = [
        ("1. Preparar", "Limpar, remover partes soltas e umedecer/saturar somente quando o procedimento exigir."),
        ("2. Fixar", "Garantir que o passante esteja firme e na posição correta antes de preencher."),
        ("3. Preencher", "Aplicar argamassa em camadas compactadas, sem deixar vazios junto ao tubo."),
        ("4. Regularizar", "Nivelar, desempenar e eliminar ondulações, ressaltos e cantos vivos."),
        ("5. Conferir", "Checar aderência, limpeza, nivelamento, cura e compatibilidade com impermeabilização."),
    ]
    x = 0.75
    for title, desc in steps:
        add_panel(slide, x, 1.65, 2.25, 4.85, fill=COLORS["white"])
        add_text(slide, title, x + 0.12, 1.95, 2.0, 0.32, size=18, bold=True, color=COLORS["blue"], align=PP_ALIGN.CENTER)
        add_text(slide, desc, x + 0.2, 2.65, 1.85, 2.25, size=13.5, color=COLORS["ink"], align=PP_ALIGN.CENTER)
        x += 2.45

    # 9
    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide)
    add_title(slide, "Checklist de liberação", "Use antes de chamar a equipe de impermeabilização.")
    checklist = [
        "[  ] O passante está fixo e sem jogo?",
        "[  ] O entorno está totalmente preenchido?",
        "[  ] Não há vazios, ninhos, trincas ou bordas quebradas?",
        "[  ] A superfície está lisa, coesa e sem partes soltas?",
        "[  ] O nível está correto, sem ondulação ou empoçamento?",
        "[  ] A área está limpa e sem pó, óleo, graxa ou detritos?",
        "[  ] A cura/tempo de liberação foi respeitada?",
        "[  ] O arremate atende ao detalhe de impermeabilização da obra?",
    ]
    add_bullets(slide, checklist, 0.9, 1.5, 6.1, 4.95, size=18, bullet=False)
    add_panel(slide, 7.55, 1.65, 4.7, 4.7, fill=COLORS["white"])
    add_text(slide, "Critério de aceite", 7.95, 2.0, 3.9, 0.35, size=22, bold=True, color=COLORS["green"], align=PP_ALIGN.CENTER)
    add_text(slide, "Passou no checklist: libera para impermeabilização.\n\nFalhou em qualquer item: corrige antes. O custo de corrigir agora é menor que reparar infiltração depois.", 8.05, 2.85, 3.75, 2.2, size=17, color=COLORS["ink"], align=PP_ALIGN.CENTER)

    # 10
    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide)
    add_title(slide, "Fixação rápida", "Perguntas para fechar o treinamento em campo.")
    add_bullets(slide, [
        "1. Por que a ondulação ao redor do passante não pode ser aceita?",
        "2. O que deve ser feito quando há vazio junto ao tubo ou à parede?",
        "3. Qual é o risco de impermeabilizar sobre poeira ou nata fraca?",
        "4. Cite três itens do checklist antes de liberar a área.",
    ], 0.9, 1.55, 6.1, 4.5, size=20, bullet=False)
    add_panel(slide, 7.45, 1.65, 4.8, 4.55, fill=COLORS["blue"])
    add_text(slide, "Mensagem-chave", 7.8, 2.05, 4.1, 0.35, size=22, bold=True, color=COLORS["white"], align=PP_ALIGN.CENTER)
    add_text(slide, "Impermeabilizante não corrige base ruim.\n\nChumbamento correto é base firme, lisa, nivelada e limpa.", 7.95, 2.9, 3.8, 1.65, size=20, bold=True, color=COLORS["white"], align=PP_ALIGN.CENTER)

    # 11
    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide)
    add_title(slide, "Referências usadas", "Conteúdo técnico resumido para treinamento interno.")
    add_bullets(slide, [
        "ABNT NBR 9575: requisitos de projeto para garantir estanqueidade e compatibilidade do sistema de impermeabilização.",
        "ABNT NBR 9574: execução de impermeabilização, incluindo preparo da base e detalhes críticos.",
        "UTFPR, Sistemas de impermeabilização: regularização uniforme, coesa, aderida, sem protuberâncias; passantes e ralos rigidamente fixados.",
        "IBI, Casos reais de preparação de superfícies: a qualidade da base é pré-requisito para o desempenho da impermeabilização.",
        "Quartzolit/Vedacit: substrato limpo, seco/íntegro conforme produto, sem pó, óleo, graxa, detritos ou irregularidades.",
    ], 0.85, 1.55, 11.6, 4.75, size=15)
    add_text(slide, "Links consultados: repositorio.utfpr.edu.br | ibibrasil.org.br | blok.com.br | quartzolit.weber | vedacit.com.br", 0.95, 6.55, 11.1, 0.24, size=9.5, color=COLORS["muted"])

    prs.save(OUT)
    print(OUT)


if __name__ == "__main__":
    missing = [str(path) for path in photos.values() if not path.exists()]
    if missing:
        raise FileNotFoundError("Fotos nao encontradas: " + "; ".join(missing))
    build()
